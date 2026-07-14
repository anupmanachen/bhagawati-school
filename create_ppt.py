#!/usr/bin/env python3
"""
Create a comprehensive PowerPoint Presentation on:
- Concept of Decentralized Planning
- Concept of Micro Level Planning
- School Improvement Plan (SIP)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(0, 102, 153)      # Dark Teal
SECONDARY_COLOR = RGBColor(0, 128, 128)    # Teal
ACCENT_COLOR = RGBColor(255, 153, 0)       # Orange
TITLE_COLOR = RGBColor(0, 51, 102)         # Navy
TEXT_COLOR = RGBColor(51, 51, 51)          # Dark Gray
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 248, 255)         # Alice Blue

def add_title_slide(prs, title, subtitle=""):
    """Add a professional title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.2))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 230, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_COLOR
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, sub_bullets=None):
    """Add a standard content slide with title and bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
        p.level = 0
        
        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = "    ◦ " + sub
                sp.font.size = Pt(16)
                sp.font.color.rgb = RGBColor(80, 80, 80)
                sp.space_after = Pt(4)
                sp.level = 1
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)
    
    return slide

def add_numbered_slide(prs, title, items):
    """Add a numbered list slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(9)
    
    return slide

def add_card_slide(prs, title, cards):
    """Add a slide with multiple cards/boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add cards
    num_cards = len(cards)
    card_width = Inches(3.9)
    card_height = Inches(2.2)
    start_x = Inches(0.5)
    start_y = Inches(1.3)
    gap = Inches(0.25)
    
    for idx, (card_title, card_content) in enumerate(cards):
        col = idx % 3
        row = idx // 3
        
        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)
        
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 250, 255)
        card.line.color.rgb = SECONDARY_COLOR
        card.line.width = Pt(1.5)
        
        # Card title
        card_title_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), card_width - Inches(0.3), Inches(0.45))
        tf = card_title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Card content
        card_content_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), card_width - Inches(0.3), card_height - Inches(0.65))
        tf = card_content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card_content
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
    
    return slide

# ============================================
# SLIDE 1: Title Slide
# ============================================
add_title_slide(
    prs,
    "विकेन्द्रित योजनाको अवधारणा\nConcept of Decentralized Planning",
    "सूक्ष्म योजनाको अवधारणा | Concept of Micro Level Planning\nविद्यालय सुधार योजना (SIP)"
)

# ============================================
# SLIDE 2: Introduction
# ============================================
add_content_slide(
    prs,
    "विकेन्द्रित योजनाको अवधारणा (Concept of Decentralized Planning)",
    [
        "विकेन्द्रित योजनाको अवधारणा ल्याटिन अमेरिकामा विकेन्द्रित शिक्षा प्रणालीबाट आएको हो।",
        "मूलभूतरूपमा त्यहाँ तीन किसिमबाट विकेन्द्रीकरणको प्रक्रिया अवलम्बन गरिएको थियो।",
        "विकेन्द्रीकरणले निर्णय अधिकार केन्द्रबाट तल्लो तहसम्म पुर्‍याउने काम गर्दछ।",
        "यसले स्थानीय आवश्यकता अनुसार शैक्षिक सेवाहरू उपलब्ध गराउँछ।",
        "नेपालमा पनि विद्यालय सुधार योजना (SIP) मार्फत विकेन्द्रीकरणको अवधारणा लागू गरिएको छ।"
    ]
)

# ============================================
# SLIDE 3: Latin America - Three Methods
# ============================================
add_section_slide(prs, "ल्याटिन अमेरिकामा विकेन्द्रीकरणका तीन विधिहरू")

# ============================================
# SLIDE 4: Three Methods Detailed
# ============================================
cards = [
    ("सूक्ष्मीकरण (Nuclearization)", "केन्द्रीय विद्यालयहरू जो बहुतै उच्च गुणस्तरका छन्, तिनीहरूमार्फत अन्य विद्यालयलाई शैक्षिक सेवाहरू उपलब्ध गराइन्छ।"),
    ("क्षेत्रीयकरण (Regionalization)", "निर्णय गर्ने अधिकार केन्द्रबाट क्षेत्रीय तहसम्म प्रत्यायोजित हुन्छ र शैक्षिक व्यवस्थापनको आधार एकाइ क्षेत्रलाई मानिन्छ।"),
    ("नगरीकरण (Municipalization)", "निर्णय गर्ने अधिकार नगरपालिकालाई हुन्छ जो वास्तविकरूपमा निर्णय अधिकार समुदायमा ल्याउन सफल साबित भएको छ।")
]
add_card_slide(prs, "ल्याटिन अमेरिकामा विकेन्द्रीकरणका तीन विधिहरू", cards)

# ============================================
# SLIDE 5: Three Stages Overview
# ============================================
add_section_slide(prs, "विकेन्द्रीकरणका तीन अवस्था (Three Stages)")

# ============================================
# SLIDE 6: Three Stages List
# ============================================
add_content_slide(
    prs,
    "विकेन्द्रीकरणको तीन अवस्था",
    [
        "विकेन्द्रीकरण (Decentralisation)",
        "प्रत्यायोजन (Delegation of Authority)",
        "निरूपण / प्रत्यारोपण (Devolution)"
    ]
)

# ============================================
# SLIDE 7: Decentralisation
# ============================================
add_content_slide(
    prs,
    "(i) विकेन्द्रीकरण (Decentralisation)",
    [
        "एउटा संगठनभित्रको विभिन्न शाखा/एकाइहरूलाई कार्यको जिम्मेवारी दिइन्छ।",
        "तर अधिकार सबै केन्द्रमा रहन्छ।",
        "केन्द्रले दिएको अधिकार कुनै पनि समयमा फिर्ता लिन सक्छ।",
        "निर्णय गर्दा माथिल्लो अधिकारीको मुख ताक्नुपर्ने हुन्छ।",
        "यो विकेन्द्रीकरणको प्रारम्भिक अवस्था हो।"
    ]
)

# ============================================
# SLIDE 8: Delegation
# ============================================
add_content_slide(
    prs,
    "(ii) अधिकार प्रत्यायोजन (Delegation of Authority)",
    [
        "माथिल्लो निकायले आफूमा निहित निर्णयको अधिकार तल्लो निकायमा प्रत्यायोजन गर्दछ।",
        "प्रत्यायोजित अधिकार प्रत्यायोजन गर्ने अधिकारीले कुनै पनि बेला झिक्न सक्छ।",
        "कर्मचारीतन्त्र माथि प्रत्यायोजन गरिन्छ।",
        "यो अवस्थामा निर्णय अधिकार केही हदसम्म तल्लो तहमा पुग्छ।",
        "तर पूर्ण स्वायत्तता प्राप्त हुँदैन।"
    ]
)

# ============================================
# SLIDE 9: Devolution
# ============================================
add_content_slide(
    prs,
    "(iii) निरूपण (Devolution) - सबैभन्दा उत्तम अवस्था",
    [
        "यो विकेन्द्रीकरणको सबैभन्दा उत्तम अवस्था हो।",
        "संस्थाहरूलाई निर्णयको अधिकार कानुनद्वारा नै व्यवस्थित गरिएको हुन्छ।",
        "अधिकारप्राप्त निकायले निर्णय गर्दा माथिल्लो निकायको मुख ताक्नु पर्दैन।",
        "कार्यान्वयन गर्ने निकाय स्वतन्त्र रूपले काम गर्न पाउँछन्।",
        "उदाहरण: गा.वि.स. लाई स्थानीय कर लगाउने अधिकार।",
        "यो पूर्ण विकेन्द्रीकरणको रूप हो।"
    ]
)

# ============================================
# SLIDE 10: Micro Level Planning
# ============================================
add_section_slide(prs, "सूक्ष्म योजनाको अवधारणा\nConcept of Micro Level Planning")

# ============================================
# SLIDE 11: Micro Planning Concept
# ============================================
add_content_slide(
    prs,
    "सूक्ष्म योजनाको अवधारणा (Micro Level Planning)",
    [
        "योजना सूक्ष्म (Micro) हो कि बृहत् (Macro), यो योजना निर्माण हुने तहमा निर्भर गर्दछ।",
        "बृहत् योजना राष्ट्रिय वा केन्द्रीय स्तरमा तयार गरिन्छ।",
        "सूक्ष्म योजना तल्लो एकाइमा तर्जुमा गरिन्छ (जिल्ला वा गाउँ स्तर)।",
        "संघीय सरकार भएको देशमा प्रान्तीय/राज्य स्तरीय सरकारलाई संविधान प्रदत्त अधिकार प्राप्त हुन्छ।",
        "एकात्मक सरकार भएको देशमा केन्द्रबाट नै शासकीय नियन्त्रण हुन्छ।",
        "भारतमा जिल्लालाई तल्लो निकाय मानिएको छ भने नेपालमा विद्यालयलाई मानिएको छ।"
    ]
)

# ============================================
# SLIDE 12: Micro vs Macro Differences
# ============================================
add_two_column_slide(
    prs,
    "सूक्ष्म र वृहत् योजनामा फरक (Micro vs Macro Planning)",
    "बृहत् योजना (Macro Planning)",
    [
        "कम बाधाहरू (Less Constraints) हुन्छन्",
        "रणनीति तथा निर्णयहरूमा विविधता हुन्छ",
        "राष्ट्रिय वा केन्द्रीय स्तरमा तयार गरिन्छ",
        "स्रोतको बाँडफाँड एवं सुविधाको विस्तारमा जोड दिन्छ",
        "उच्च स्तरको योजना एकाइबाट निर्देशित"
    ],
    "सूक्ष्म योजना (Micro Planning)",
    [
        "आधार एवं सीमा (Parameters & Limitations) हुन्छन्",
        "माथिल्लो स्तरको योजना एकाइबाट निर्देशित",
        "तल्लो निकायलाई अधिक स्वतन्त्रता",
        "स्थानीय आवश्यकता र चाहनालाई अङ्गीकार",
        "जनसहभागिता एवं लागत सहभागिता"
    ]
)

# ============================================
# SLIDE 13: Objectives of Micro Planning
# ============================================
add_numbered_slide(
    prs,
    "सूक्ष्म योजनाको उद्देश्य (Objectives of Micro-planning)",
    [
        "योजनालाई स्थानीय आवश्यकतामा आधारित, अत्याधिक सहभागितामूलक बनाउन स्थानीय समुदायलाई परिचालन गरिन्छ।",
        "विद्यालय र समुदायबीच सुमधुर सम्बन्ध तथा अन्तरक्रिया विकसित गरी विद्यालय तथा शिक्षकहरूलाई निरन्तररूपमा प्राप्त हुने उत्तम सहयोग प्रणाली (Support System) को स्थापना हुन्छ।",
        "तोकिएका उमेर समूहका सम्पूर्ण बालबालिकाहरूलाई विद्यालय वा वैकल्पिक विद्यालय भित्र प्रवेश गर्ने सुनिश्चितता प्राप्त हुन्छ।",
        "विद्यालयको कार्य नियमित र सक्षमरूपमा सञ्चालनको सुनिश्चितता प्रदान गरिन्छ।",
        "सूक्ष्म योजना एकाइले विद्यालयको सेवा क्षेत्रका गाउँहरूसँग राम्रो समझदारी कायम राख्नुपर्दछ।",
        "जनमुखी, सहभागितामूलक, स्थानीय आवश्यकतामा आधारित स्रोत सदुपयोग एवं संस्थागत सक्षमता सुधारमा अवलम्बित हुन्छ।"
    ]
)

# ============================================
# SLIDE 14: SIP Introduction
# ============================================
add_section_slide(prs, "विद्यालय सुधार योजना\nSchool Improvement Plan (SIP)")

# ============================================
# SLIDE 15: SIP Importance
# ============================================
add_content_slide(
    prs,
    "विद्यालय सुधार योजना (SIP) - परिचय",
    [
        "विद्यालयहरूको व्यवस्थापकीय कुशलतामा अभिवृद्धि गरी उपलब्ध स्रोत र साधनहरूको सर्वोत्तम परिचालन गरेर उच्चतम प्रतिफल सुनिश्चित गर्नुपर्ने यथार्थता आजको सबैभन्दा ठूलो चासोको विषय हो।",
        "विद्यालयले आफूलाई आफैं सुधार नगरेसम्म केन्द्रीय तहबाट गरिएका सुधारका प्रयत्नले खासै प्रतिफल दिन सक्दैन।",
        "विद्यालय सुधार योजना भनेको विद्यालयलाई बदलिँदो परिस्थितिको मागअनुसार विद्यार्थीहरूको पहुँच एवम् सिकाइ उपलब्धिमा सुधार ल्याउन तयार गरिने सूक्ष्म योजना हो।",
        "यसले विद्यालयको भौतिक निर्माण कार्यमा मात्र होइन, सबै पक्षहरूको सर्वाङ्गीण विकासद्वारा विद्यार्थीको उच्च सिकाइ उपलब्धि सुनिश्चित गर्दछ।"
    ]
)

# ============================================
# SLIDE 16: Characteristics of Improved Schools - Part 1
# ============================================
add_content_slide(
    prs,
    "सुधारिएका विद्यालयका विशेषताहरू (Part 1)",
    [
        "स्पष्ट दृष्टिकोण (Clear Vision): सुदूर भविष्यको परिकल्पनालाई साकार पार्न प्रयास गर्दछ।",
        "स्पष्ट उद्देश्य (Clear Objectives): चरम उद्देश्य निर्धारण गरेको हुन्छ र सम्बन्धितहरूबीच सहमति कायम गर्दछ।",
        "रणनीतिहरू (Strategies): उद्देश्य प्राप्तिका लागि आवश्यक उपयुक्त रणनीति पहिले नै निर्धारण गर्दछ।",
        "कार्यमूलक जिम्मेवारी र कार्यक्षेत्रको विभाजन: गतिविधिहरूबाट प्रभावित सबै पक्षको जिम्मेवारी स्पष्टरूपले विभाजन गरिएको हुन्छ।",
        "सरोकारवालाहरूको सक्रिय सहभागिता: निर्णय प्रक्रियामा विद्यार्थी, शिक्षक, अभिभावक तथा समुदायले सक्रिय सहभागिता जनाउँछन्।"
    ]
)

# ============================================
# SLIDE 17: Characteristics of Improved Schools - Part 2
# ============================================
add_content_slide(
    prs,
    "सुधारिएका विद्यालयका विशेषताहरू (Part 2)",
    [
        "सक्रिय सिकाइ (Active Learning): शिक्षकको भूमिका विद्यार्थीहरूको सहयोगीको रूपमा रहन्छ।",
        "स्तर निर्धारण (Standard Fixation): क्रियाकलापहरूको व्यवहार र कार्यसम्पादनको स्पष्ट स्तर निर्धारण गरिएको हुन्छ।",
        "स्रोत र साधनको परिचालन (Resource Mobilization): आवश्यक स्रोत र साधन कसरी उपलब्ध गर्ने र अधिकतम उपयोग गर्ने भन्ने बारे स्पष्ट हुन्छ।",
        "विद्यालय व्यवस्थापन (School Management): शैक्षिक योजना तयार गर्ने, कार्यक्रम तर्जुमा गर्ने र समन्वित कार्यान्वयन गर्ने।",
        "शैक्षिक व्यवस्थापन सूचना प्रणाली (EMIS): शिक्षासँग सम्बन्धित सबै सूचना तथा तथ्याङ्कहरूलाई व्यवस्थित रूपमा सङ्कलन र व्यवस्थापन गर्दछ।"
    ]
)

# ============================================
# SLIDE 18: Planning for School Improvement
# ============================================
add_content_slide(
    prs,
    "विद्यालय सुधारको लागि योजना निर्माण",
    [
        "विद्यालयका उपभोक्ताहरू (अभिभावक, शिक्षक, विद्यार्थी तथा समुदाय) सबैसँगै बसी आवश्यक रणनीति निर्माणमा साझा प्रयास अगाडि बढाउन सकिन्छ।",
        "विद्यालयको मौजुदा अवस्थालाई केलाएर मुख्य चुनौतीहरू पहिचान गर्नु पर्छ।",
        "विद्यालयलाई उद्देश्यअनुरुप अगाडि बढ्न के-कस्ता बाधाहरू आएका छन्? तिनको निराकरण गर्ने उपायहरू के-के हुन्?",
        "यसबाट योजना निर्माण र कार्यान्वयन गर्न आवश्यक रणनीति निर्माणका लागि साझा प्रयास अगाडि बढाउन सकिन्छ।",
        "विद्यालय सुधार योजनाले विद्यालयको भौतिक सुधारदेखि लिएर शिक्षक र विद्यार्थीको नियमितता, उनीहरूको सिकाइ उपलब्धि अभिवृद्धि, विद्यालयको आन्तरिक सक्षमता लगायत सबै पक्षको सुधारलाई लक्षित गर्नुपर्दछ।"
    ]
)

# ============================================
# SLIDE 19: SIP Process
# ============================================
add_content_slide(
    prs,
    "विद्यालय सुधार योजना निर्माण प्रक्रिया",
    [
        "सरोकारवालाहरूका सबै पक्षहरू एकै साथ बसी आफ्ना शैक्षिक आवश्यकता पहिचान गर्ने।",
        "आवश्यकताहरूको प्राथमिकता निर्धारण गर्ने।",
        "उद्देश्य निर्धारण गर्ने।",
        "ती उद्देश्यप्राप्तिका लागि आफैले योजना निर्माण गर्ने।",
        "योजनाको कार्यान्वयन र अनुगमन एवम् मूल्याङ्कन गर्नुपर्छ।",
        "यसबाट सरोकारवालाहरूमा स्वामित्वको भावना जागृत भएर कार्यक्रममा दीगोपनाको विकास हुन्छ।",
        "विद्यालय प्रणालीको प्रभावकारिता सुनिश्चित हुन्छ।"
    ]
)

# ============================================
# SLIDE 20: Steps of SIP
# ============================================
add_numbered_slide(
    prs,
    "विद्यालय सुधार योजनाका चरणहरू (Steps of SIP)",
    [
        "विद्यालय सेवा क्षेत्रको सर्वेक्षण गरी विद्यमान शैक्षिक अवस्थाको आवश्यकताको पहिचान एवम् विश्लेषण।",
        "विद्यालयका मुख्य चुनौतीहरू र कमी-कमजोरीहरूको पहिचान।",
        "उद्देश्य तथा लक्ष्य निर्धारण (पहुँच तथा सहभागिता, सिकाइ उपलब्धि, व्यवस्थापन कुशलता)।",
        "आवश्यक रणनीति तथा कार्यनीति निर्माण।",
        "स्रोत तथा साधनको आवश्यकता आँकलन र व्यवस्था।",
        "कार्यक्रम तर्जुमा, कार्यान्वयन योजना निर्माण।",
        "अनुगमन, मूल्याङ्कन र समीक्षा प्रक्रिया।"
    ]
)

# ============================================
# SLIDE 21: Conclusion
# ============================================
add_content_slide(
    prs,
    "निष्कर्ष (Conclusion)",
    [
        "विकेन्द्रीकरणले शैक्षिक निर्णय अधिकार तल्लो तहमा ल्याउँछ र स्थानीय आवश्यकता अनुसार काम गर्न मद्दत गर्दछ।",
        "सूक्ष्म योजना (Micro Planning) ले स्थानीय समुदायको सहभागितालाई सुनिश्चित गर्दछ।",
        "विद्यालय सुधार योजना (SIP) विद्यालयको सर्वाङ्गीण विकास र उच्च सिकाइ उपलब्धिका लागि महत्वपूर्ण उपकरण हो।",
        "सुधारिएका विद्यालयहरूमा स्पष्ट दृष्टिकोण, उद्देश्य, सक्रिय सहभागिता र स्रोत परिचालन जस्ता विशेषताहरू हुन्छन्।",
        "SIP ले विद्यालयलाई आत्मनिर्भर, सहभागितामूलक र प्रभावकारी बनाउँछ।"
    ]
)

# ============================================
# SLIDE 22: Thank You
# ============================================
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Background
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_COLOR
shape.line.fill.background()

# Thank you text
thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = thank_box.text_frame
p = tf.paragraphs[0]
p.text = "धन्यवाद"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(200, 230, 255)
p.alignment = PP_ALIGN.CENTER

# Save the presentation
output_path = "/home/user/bhagawati-school/Decentralized_Planning_Micro_Planning_SIP.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to: {output_path}")
print(f"Total slides: {len(prs.slides)}")